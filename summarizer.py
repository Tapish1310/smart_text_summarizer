import os
import textwrap
import fitz
import evaluate
import nltk
import concurrent.futures
from nltk.tokenize import sent_tokenize
from docx import Document
from pptx import Presentation
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from reportlab.lib.pagesizes import A4
from reportlab.pdfgen import canvas

nltk.download('punkt')

MODELS = {
    "1": "facebook/bart-large-cnn",
    "2": "t5-base",
    "3": "google/pegasus-xsum"
}

def load_model(model_choice):
    model_name = MODELS.get(model_choice, MODELS["1"])
    tokenizer = AutoTokenizer.from_pretrained(model_name)
    model = AutoModelForSeq2SeqLM.from_pretrained(model_name)
    return tokenizer, model, model_name

def apply_style_prompt(text, style, model_name):
    if model_name.startswith("t5"):
        if style == "2":
            return "summarize in bullet points: " + text
        elif style == "3":
            return "summarize as headlines: " + text
        else:
            return "summarize: " + text
    else:
        return text

def post_process(summary, style):
    if style == "2":
        return "\n• " + "\n• ".join(sent_tokenize(summary))
    elif style == "3":
        return "\n- " + "\n- ".join(sent_tokenize(summary))
    else:
        return summary

def summarize_chunk(text, tokenizer, model, max_len, min_len):
    inputs = tokenizer(
        text,
        return_tensors="pt",
        truncation=True,
        padding="max_length",
        max_length=1024
    )
    if inputs["input_ids"].shape[-1] > 1024:
        print("⚠️ Chunk skipped: too long for model.")
        return ""
    summary_ids = model.generate(
        inputs["input_ids"],
        max_length=max_len,
        min_length=min_len,
        num_beams=4,
        early_stopping=True
    )
    return tokenizer.decode(summary_ids[0], skip_special_tokens=True)

def sentence_chunking(text, max_words=500, overlap=50):
    sentences = sent_tokenize(text)
    chunks, chunk = [], []
    word_count = 0

    for sentence in sentences:
        words = sentence.split()
        if word_count + len(words) > max_words:
            chunks.append(" ".join(chunk))
            chunk = chunk[-overlap:]
            word_count = sum(len(s.split()) for s in chunk)

        chunk.append(sentence)
        word_count += len(words)

    if chunk:
        chunks.append(" ".join(chunk))
    return chunks

def summarize_all_chunks(chunks, tokenizer, model, max_len, min_len):
    with concurrent.futures.ThreadPoolExecutor() as executor:
        futures = [
            executor.submit(summarize_chunk, chunk, tokenizer, model, max_len, min_len)
            for chunk in chunks
        ]
        return [f.result() for f in concurrent.futures.as_completed(futures)]

def extract_text_from_file(file_path):
    if not os.path.exists(file_path):
        print("❌ File not found.")
        return None
    try:
        if file_path.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        elif file_path.endswith(".pdf"):
            pdf_doc = fitz.open(file_path)
            return "".join([page.get_text() for page in pdf_doc])
        elif file_path.endswith(".docx"):
            doc = Document(file_path)
            return "\n".join([p.text for p in doc.paragraphs])
        elif file_path.endswith(".pptx"):
            prs = Presentation(file_path)
            return "\n".join([
                shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
            ])
        else:
            print("❌ Unsupported file format.")
            return None
    except Exception as e:
        print(f"❌ Error: {e}")
        return None

def evaluate_summary(summary, reference):
    rouge = evaluate.load("rouge")
    return rouge.compute(predictions=[summary], references=[reference])

def save_summary_as_pdf(text, filename="summary_output.pdf"):
    c = canvas.Canvas(filename, pagesize=A4)
    width, height = A4
    x_margin, y_margin = 50, 800
    lines = textwrap.wrap(text, width=100)

    for line in lines:
        if y_margin < 50:
            c.showPage()
            y_margin = 800
        c.drawString(x_margin, y_margin, line)
        y_margin -= 15

    c.save()
    print(f"📄 PDF saved as: {filename}")
